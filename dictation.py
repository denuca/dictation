import re
import json
import os
import pptx
import sys
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Cm
from PIL import ImageFont

# Load styles from JSON file
with open('styles.json', 'r') as f:
    styles = json.load(f)

# Function to create a PowerPoint presentation from a text file
def create_ppt_from_text(file_path, output_path='output.pptx'):
    # Load the text file
    with open(file_path, 'r') as file:
        lines = file.readlines()

    # Create a presentation object
    prs = pptx.Presentation()

    # Initialize slide number and word ID
    slide_number = 1
    word_id = 1

    # Create a new slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = lines[0].strip()  # First line as title
    content_placeholder = slide.placeholders[1]

    # Function to calculate text size using Pillow
    def calculate_text_size(text, font_path, font_size):
        font = ImageFont.truetype(font_path, font_size)
        bbox = font.getbbox(text)
        width = bbox[2] - bbox[0]
        height = bbox[3] - bbox[1]
        # Add some padding
        width += 10
        height += 10
        return width, height

    # Function to add a text box with a border and adjusted margins
    def add_text_box(slide, text, left, top, word_id, max_height, content_type):
        # Calculate text size
        text_width, text_height = calculate_text_size(text, styles[content_type]['font_path'], styles[content_type]['font_size'])
        text_width = Inches(text_width / 72)  # Convert points to inches

        # Add the text box
        txBox = slide.shapes.add_textbox(left, top, text_width, max_height)
        tf = txBox.text_frame
        tf.text = text
        for paragraph in tf.paragraphs:
            paragraph.font.name = styles[content_type]['font_name']
            paragraph.font.size = Pt(styles[content_type]['font_size'])
            paragraph.font.color.rgb = RGBColor(*styles[content_type]['font_color'])

        txBox.line.color.rgb = RGBColor(*styles[content_type]['border_color'])
        txBox.line.width = Pt(styles[content_type]['border_width'])

        # Adjust text box margins
        tf.margin_left = Pt(styles[content_type]['margin_left'])
        tf.margin_right = Pt(styles[content_type]['margin_right'])
        tf.margin_top = Pt(styles[content_type]['margin_top'])
        tf.margin_bottom = Pt(styles[content_type]['margin_bottom'])

        # Disable AutoFit
        tf.word_wrap = False
        #tf.auto_size = True

        # Add the word ID above the text box if required
        if styles[content_type]['display_id']:
            id_box = slide.shapes.add_textbox(left, top - Inches(0.3), text_width, Inches(0.3))
            id_frame = id_box.text_frame
            id_frame.text = str(word_id)
            id_frame.paragraphs[0].font.size = Pt(styles['id']['font_size'])
            id_frame.paragraphs[0].font.color.rgb = RGBColor(*styles['id']['font_color'])
            id_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Process each line in the text file
    max_height = 0
    for line in lines[1:]:
        words = re.findall(r'\w+|[^\w\s]', line, re.UNICODE)
        for word in words:
            content_type = 'word' if re.match(r'\w+', word) else 'punctuation'
            _, text_height = calculate_text_size(word, styles[content_type]['font_path'], styles[content_type]['font_size'])
            max_height = max(max_height, Inches(text_height / 72))

    for line in lines[1:]:
        words = re.findall(r'\w+|[^\w\s]', line, re.UNICODE)
        left = content_placeholder.left
        top = content_placeholder.top
        for word in words:
            content_type = 'word' if re.match(r'\w+', word) else 'punctuation'
            add_text_box(slide, word, left, top + Inches(0.3), word_id, max_height, content_type)
            text_width, _ = calculate_text_size(word, styles[content_type]['font_path'], styles[content_type]['font_size'])
            left += Inches(text_width / 72) + Inches(0.2)  # Adjust spacing between text boxes
            if styles[content_type]['count_id']:
                word_id += 1

            # Check if the next text box will fit in the current line
            if left + Inches(1.5) > content_placeholder.left + content_placeholder.width:
                left = content_placeholder.left
                top += max_height + Inches(0.6)  # Move to the next line

                # Check if the next text box will fit in the current slide
                if top + max_height + Inches(0.6) > content_placeholder.top + content_placeholder.height:
                    slide_number += 1
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    title = slide.shapes.title
                    title.text = lines[0].strip()
                    content_placeholder = slide.placeholders[1]
                    top = content_placeholder.top

    # Delete the content placeholders from all slides as they were used for dimensions only.
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.idx == 1:
                sp = shape._element
                sp.getparent().remove(sp)

    # Add headers and footers
    for i, slide in enumerate(prs.slides):
        # Add header
        header = slide.shapes.add_textbox(Cm(1), Cm(0.5), prs.slide_width - Cm(2), Cm(1))
        header_frame = header.text_frame
        header_frame.text = styles['header']['text']
        header_frame.paragraphs[0].font.size = Pt(styles['header']['font_size'])
        header_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Add footer with page number
        footer = slide.shapes.add_textbox(Cm(1), prs.slide_height - Cm(1.5), prs.slide_width - Cm(2), Cm(1))
        footer_frame = footer.text_frame
        footer_frame.text = f"Page {i + 1} of {len(prs.slides)}"
        footer_frame.paragraphs[0].font.size = Pt(styles['footer']['font_size'])
        footer_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    # Save the presentation
    prs.save(output_path)
    pass

if __name__ == '__main__':
    if len(sys.argv) > 1:
        print(f"Argument is: '{sys.argv[1]}'.")
        create_ppt_from_text(f"{sys.argv[1]}")
    else:
        print("Please provide the path to the text file as an argument.")
